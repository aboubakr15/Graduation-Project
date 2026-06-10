import os
import logging
from typing import Optional, List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy

logger = logging.getLogger(__name__)

# Try to import MSO_SHAPE
try:
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    from pptx.util import MSO_SHAPE


class PresentationMaker:
    """
    Handles the creation of professional PowerPoint presentations.
    Supports multiple slide types: cover, content, two_column, closing.
    """

    def __init__(self, output_dir="presentations"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

        # ── Brand Colors ──────────────────────────────────────────────
        self.COLOR_BG         = RGBColor(0x1E, 0x27, 0x61)   # Deep Navy
        self.COLOR_BG_ALT     = RGBColor(0x14, 0x1A, 0x45)   # Darker Navy (cover/closing)
        self.COLOR_TITLE      = RGBColor(0xFF, 0xFF, 0xFF)   # White
        self.COLOR_BODY       = RGBColor(0xCA, 0xDC, 0xFC)   # Ice Blue
        self.COLOR_ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # Cyan Accent
        self.COLOR_ACCENT2    = RGBColor(0xFF, 0xC3, 0x00)   # Gold (closing)
        self.COLOR_MUTED      = RGBColor(0x94, 0xA3, 0xB8)   # Slate (slide numbers, captions)
        self.COLOR_CARD_BG    = RGBColor(0x25, 0x32, 0x75)   # Card background (slightly lighter)

        # ── Fonts ─────────────────────────────────────────────────────
        self.FONT_TITLE  = "Calibri"
        self.FONT_BODY   = "Calibri Light"

    # ──────────────────────────────────────────────────────────────────
    # Public API
    # ──────────────────────────────────────────────────────────────────

    def create_presentation(
        self,
        slides_data: List[Dict],
        image_paths: List[str] = None,
        filename: str = "presentation.pptx",
        slide_image_map: Dict[int, str] = None,
    ) -> Optional[str]:
        """
        Creates a professionally styled .pptx file.

        Each item in slides_data should have:
            - type:    "cover" | "content" | "two_column" | "closing"  (default: "content")
            - title:   str
            - content: list[str]  – bullet points
            - notes:   str        – speaker notes (optional)

        image_paths: optional list of local image file paths provided by the user.
            - If None or empty  → no images are used anywhere, period.
            - If provided       → images are distributed across non-cover/non-closing slides
                                  in order. Slides without a matching image render as
                                  plain content slides (no image inserted).

        slide_image_map: optional dict {slide_index: image_path} for PRECISE per-slide
            image assignment (overrides image_paths ordering). Used for code-image slides.
        """
        try:
            prs = Presentation()
            # Widescreen 16:9
            prs.slide_width  = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Build the final img_map: slide index → image path
            img_map: Dict[int, str] = {}

            if slide_image_map:
                # Precise per-slide assignment (code images)
                for idx, p in slide_image_map.items():
                    if p and os.path.exists(p):
                        img_map[idx] = p
            elif image_paths:
                # Legacy: distribute valid images to content slides in order
                valid_images: List[str] = [
                    p for p in image_paths if p and os.path.exists(p)
                ]
                if valid_images:
                    image_iter = iter(valid_images)
                    for idx, slide_info in enumerate(slides_data):
                        stype = slide_info.get("type", "content").lower()
                        if stype not in ("cover", "closing"):
                            img = next(image_iter, None)
                            if img:
                                img_map[idx] = img

            for idx, slide_info in enumerate(slides_data):
                slide_type = slide_info.get("type", "content").lower()
                img_path   = img_map.get(idx)   # None if no image assigned

                if slide_type == "cover":
                    self._add_cover_slide(prs, slide_info, idx + 1, len(slides_data))
                elif slide_type == "closing":
                    self._add_closing_slide(prs, slide_info, idx + 1, len(slides_data))
                elif img_path:
                    # Image available → two-column layout (text left, image right)
                    self._add_two_column_slide(prs, slide_info, img_path, idx + 1, len(slides_data))
                else:
                    # Plain content slide
                    self._add_content_slide(prs, slide_info, idx + 1, len(slides_data))

            output_path = self._unique_path(os.path.join(self.output_dir, filename))
            prs.save(output_path)
            logger.info(f"Presentation saved to: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Error creating presentation: {e}", exc_info=True)
            return None

    # ──────────────────────────────────────────────────────────────────
    # Markdown presentation support
    # ──────────────────────────────────────────────────────────────────

    @staticmethod
    def parse_markdown_slides(markdown_str: str) -> List[Dict]:
        """
        Parse a Markdown presentation string into a list of slide dicts
        compatible with create_presentation().

        Expected format::

            # Slide Title
            - Bullet 1
            - Bullet 2

            <!-- slide -->

            # Next Slide Title
            - Bullet 1

        Rules:
        - Slides are separated by ``<!-- slide -->`` (case-insensitive).
        - First block  → type = "cover"
        - Last block   → type = "closing"  (auto-detected or forced)
        - Other blocks → type = "content"
        - Lines starting with ``#``   → slide title (first heading wins)
        - Lines starting with ``-``/``*`` → bullet points
        - Fenced code blocks (```...```) → stored in "code_block" field
        - Empty / other lines → ignored
        """
        import re as _re

        raw_blocks = _re.split(r'<!--\s*slide\s*-->', markdown_str, flags=_re.IGNORECASE)
        slides = []

        for raw in raw_blocks:
            block = raw.strip()
            if not block:
                continue

            title = ""
            bullets = []
            code_block = ""

            lines = block.splitlines()
            i = 0
            while i < len(lines):
                line = lines[i]
                stripped = line.strip()
                if not stripped:
                    i += 1
                    continue
                # Detect fenced code blocks
                if stripped.startswith('```'):
                    # Collect all lines until closing fence
                    lang = stripped[3:].strip()  # e.g. 'python'
                    code_lines = []
                    i += 1
                    while i < len(lines) and not lines[i].strip().startswith('```'):
                        code_lines.append(lines[i])
                        i += 1
                    # i is now on the closing ``` line
                    code_block = "\n".join(code_lines).strip()
                    i += 1
                    continue
                if stripped.startswith('#'):
                    if not title:
                        title = stripped.lstrip('#').strip()
                elif stripped.startswith('-') or stripped.startswith('*'):
                    text = stripped.lstrip('-*').strip()
                    if text:
                        bullets.append(text)
                i += 1

            if not title and not bullets and not code_block:
                continue

            slide_dict = {
                "type": "content",
                "title": title or "Slide",
                "content": bullets,
                "notes": "",
            }
            if code_block:
                slide_dict["code_block"] = code_block

            slides.append(slide_dict)

        if not slides:
            return slides

        # Assign cover/closing types
        slides[0]["type"] = "cover"
        if len(slides) > 1:
            last = slides[-1]
            if last["title"].lower() in ("thank you", "thanks", "closing", "questions"):
                last["type"] = "closing"
            else:
                slides.append({
                    "type": "closing",
                    "title": "Thank You",
                    "content": ["Questions & Discussion"],
                    "notes": "",
                })

        return slides

    def generate_code_images_from_markdown(self, markdown_str: str) -> List[str]:
        """
        Scan the markdown for code-slide blocks and render each code block
        as a PNG image (dark background, light monospace text).

        Returns:
            List of absolute paths to the generated temporary PNG images.
            Returns [] if Pillow is not available or no code blocks found.
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            logger.warning("Pillow not installed — code images cannot be generated. Install with: pip install Pillow")
            return []

        import re as _re
        import tempfile

        CODE_SLIDE_KEYWORDS = [
            "code", "implementation", "example", "algorithm",
            "snippet", "solution", "walkthrough",
        ]

        slides = self.parse_markdown_slides(markdown_str)
        image_paths: List[str] = []

        # We only generate images for slides that actually have a code_block.
        for slide in slides:
            code = slide.get("code_block", "").strip()
            title_lower = slide.get("title", "").lower()
            is_code_slide = code or any(kw in title_lower for kw in CODE_SLIDE_KEYWORDS)
            if not is_code_slide or not code:
                image_paths.append("")  # placeholder for non-code slides
                continue

            img_path = self._render_code_image(code)
            image_paths.append(img_path if img_path else "")

        return image_paths

    def _render_code_image(
        self,
        code: str,
        width: int = 1100,
        bg_color: tuple = (18, 18, 18),
        text_color: tuple = (212, 212, 212),
        keyword_color: tuple = (86, 156, 214),
        string_color: tuple = (206, 145, 120),
        comment_color: tuple = (106, 153, 85),
        font_size: int = 20,
    ) -> Optional[str]:
        """Render a code string as a styled PNG image and return the file path."""
        try:
            from PIL import Image, ImageDraw, ImageFont
            import tempfile, os as _os

            # ── Font selection ────────────────────────────────────────────
            # Try common monospace fonts; fall back to Pillow default.
            mono_fonts = [
                "Consolas", "DejaVuSansMono", "DejaVu Sans Mono",
                "Courier New", "Courier", "Liberation Mono",
            ]
            font = None
            for fname in mono_fonts:
                try:
                    font = ImageFont.truetype(fname, font_size)
                    break
                except Exception:
                    try:
                        font = ImageFont.truetype(fname + ".ttf", font_size)
                        break
                    except Exception:
                        continue
            if font is None:
                try:
                    font = ImageFont.load_default(size=font_size)
                except Exception:
                    font = ImageFont.load_default()

            # ── Measure text dimensions ───────────────────────────────────
            padding = 24
            line_height = font_size + 8
            lines = code.splitlines()
            height = max(line_height * (len(lines) + 2) + padding * 2, 200)

            # ── Draw image ────────────────────────────────────────────────
            img = Image.new("RGB", (width, height), color=bg_color)
            draw = ImageDraw.Draw(img)

            # Header bar
            header_h = 32
            draw.rectangle([(0, 0), (width, header_h)], fill=(40, 40, 40))
            for cx, col in [(16, (255, 95, 86)), (40, (255, 189, 46)), (64, (39, 201, 63))]:
                draw.ellipse([(cx - 7, header_h // 2 - 7), (cx + 7, header_h // 2 + 7)], fill=col)
            draw.text((88, header_h // 2 - font_size // 2), "code", font=font, fill=(150, 150, 150))

            # Code lines
            y = header_h + padding

            import re as _re
            KEYWORDS = {
                "python": [
                    r"\b(def|class|return|if|elif|else|for|while|import|from|as|"
                    r"with|try|except|finally|pass|break|continue|yield|lambda|"
                    r"True|False|None|and|or|not|in|is)\b"
                ],
            }

            for line in lines:
                # Simple tokeniser: draw keyword-coloured spans, then rest
                x = padding
                remaining = line

                # Draw line number (muted)
                draw.text((x, y), " ", font=font, fill=(80, 80, 80))

                # Detect comments
                comment_pos = line.find("#")
                code_part = line if comment_pos == -1 else line[:comment_pos]
                comment_part = "" if comment_pos == -1 else line[comment_pos:]

                # Draw the non-comment part word by word
                draw.text((x, y), code_part, font=font, fill=text_color)

                # Draw comment in comment colour
                if comment_part:
                    try:
                        code_bbox = draw.textbbox((x, y), code_part, font=font)
                        comment_x = code_bbox[2]
                    except AttributeError:
                        comment_x = x + len(code_part) * (font_size // 2)
                    draw.text((comment_x, y), comment_part, font=font, fill=comment_color)

                y += line_height

            # Save to temp file
            tmp = tempfile.NamedTemporaryFile(
                suffix=".png", delete=False, prefix="code_slide_"
            )
            img.save(tmp.name, "PNG")
            tmp.close()
            logger.info(f"Code image saved: {tmp.name}")
            return tmp.name

        except Exception as e:
            logger.error(f"Failed to render code image: {e}")
            return None

    def create_from_markdown(
        self,
        markdown_str: str,
        image_paths: List[str] = None,
        filename: str = "presentation.pptx",
        per_slide_images: List[str] = None,
    ) -> Optional[str]:
        """
        Convenience wrapper: parse Markdown string and produce a .pptx file.

        Args:
            markdown_str:     Full Markdown deck (``# Title / - bullet / <!-- slide -->``).
            image_paths:      Optional list of local image paths to distribute in order.
            filename:         Output filename for the .pptx.
            per_slide_images: Optional slide-parallel list of image paths (one entry per
                              parsed slide, empty string = no image for that slide).
                              When provided, images are placed on their matching slides
                              precisely (used for code images).

        Returns:
            Absolute path to the saved .pptx, or ``None`` on error.
        """
        slides_data = self.parse_markdown_slides(markdown_str)
        if not slides_data:
            logger.error("Markdown parsing returned no slides — aborting PPTX creation.")
            return None
        logger.info(f"Parsed {len(slides_data)} slide(s) from Markdown.")

        # Build a precise slide_image_map from per_slide_images if provided
        slide_image_map: Optional[Dict[int, str]] = None
        if per_slide_images:
            slide_image_map = {}
            for idx, p in enumerate(per_slide_images):
                if p and os.path.exists(p):
                    slide_image_map[idx] = p

        return self.create_presentation(
            slides_data,
            image_paths=image_paths,
            filename=filename,
            slide_image_map=slide_image_map,
        )

    # ──────────────────────────────────────────────────────────────────
    # Slide builders
    # ──────────────────────────────────────────────────────────────────

    def _add_cover_slide(self, prs, slide_info, slide_num, total):
        """Full-bleed dark cover with large title and subtitle."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        W, H = prs.slide_width, prs.slide_height

        # Background
        self._set_bg(slide, self.COLOR_BG_ALT)

        # Left accent bar (thick)
        self._add_rect(slide, 0, 0, Inches(0.18), H, self.COLOR_ACCENT)

        # Decorative bottom bar
        self._add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), self.COLOR_ACCENT)

        # Title
        title_text = slide_info.get("title", "Presentation")
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.8), Inches(12.0), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(52)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_TITLE

        # Subtitle (first bullet treated as subtitle)
        content = self._coerce_list(slide_info.get("content", []))
        subtitle = content[0] if content else ""
        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(4.0), Inches(10.0), Inches(0.9))
            tf2 = tb2.text_frame
            p2  = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            r2  = p2.add_run()
            r2.text = subtitle
            r2.font.name  = self.FONT_BODY
            r2.font.size  = Pt(22)
            r2.font.color.rgb = self.COLOR_BODY

        # Divider line between title and subtitle
        self._add_rect(slide, Inches(0.55), Inches(3.85), Inches(5.0), Inches(0.04), self.COLOR_ACCENT)

        # Slide counter
        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    def _add_content_slide(self, prs, slide_info, slide_num, total):
        """Standard content slide with left accent bar and bullet list."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        W, H = prs.slide_width, prs.slide_height

        self._set_bg(slide, self.COLOR_BG)

        # Top accent bar
        self._add_rect(slide, 0, 0, W, Inches(0.08), self.COLOR_ACCENT)

        # Left accent bar
        self._add_rect(slide, 0, Inches(0.08), Inches(0.08), H - Inches(0.16), self.COLOR_ACCENT)

        # Title — centered across full width
        title_text = slide_info.get("title", "")
        tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(36)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_TITLE

        # Thin divider — centered under title, width scales with title length
        char_width_inches = 0.22
        divider_w = min(max(len(title_text) * char_width_inches, 1.5), 10.0)
        divider_left = (prs.slide_width.inches - divider_w) / 2
        self._add_rect(slide, Inches(divider_left), Inches(1.05), Inches(divider_w), Inches(0.02), self.COLOR_ACCENT)

        # Content bullets
        content = self._coerce_list(slide_info.get("content", []))
        if content:
            tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.0), Inches(6.0))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for i, point in enumerate(content):
                para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                para.space_before = Pt(6)
                para.space_after  = Pt(4)
                # Bullet dot
                run_dot = para.add_run()
                run_dot.text = "◆  "
                run_dot.font.name  = self.FONT_TITLE
                run_dot.font.size  = Pt(10)
                run_dot.font.color.rgb = self.COLOR_ACCENT
                # Bullet text
                run_txt = para.add_run()
                run_txt.text = point
                run_txt.font.name  = self.FONT_BODY
                run_txt.font.size  = Pt(18)
                run_txt.font.color.rgb = self.COLOR_BODY

        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    def _add_two_column_slide(self, prs, slide_info, img_path, slide_num, total):
        """Two-column slide: text left, image right.
        This method is only called when img_path is a valid, existing file."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = prs.slide_width, prs.slide_height

        self._set_bg(slide, self.COLOR_BG)

        # Top accent bar
        self._add_rect(slide, 0, 0, W, Inches(0.08), self.COLOR_ACCENT)

        # Card background for right column
        self._add_rect(slide, Inches(7.2), Inches(0.08), Inches(6.133), H - Inches(0.08), self.COLOR_CARD_BG)

        # Title
        title_text = slide_info.get("title", "")
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6.6), Inches(1.0))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(34)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_TITLE

        char_width_inches = 0.21
        divider_w = min(max(len(title_text) * char_width_inches, 1.5), 6.0)
        self._add_rect(slide, Inches(0.3), Inches(1.25), Inches(divider_w), Inches(0.02), self.COLOR_ACCENT)

        # Bullets (left column)
        content = self._coerce_list(slide_info.get("content", []))
        if content:
            tb2 = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(6.7), Inches(5.8))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for i, point in enumerate(content):
                para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                para.space_before = Pt(7)
                para.space_after  = Pt(5)
                run_dot = para.add_run()
                run_dot.text = "◆  "
                run_dot.font.size  = Pt(9)
                run_dot.font.color.rgb = self.COLOR_ACCENT
                run_txt = para.add_run()
                run_txt.text = point
                run_txt.font.name  = self.FONT_BODY
                run_txt.font.size  = Pt(16)
                run_txt.font.color.rgb = self.COLOR_BODY

        # Insert user-provided image on the right column
        try:
            slide.shapes.add_picture(
                img_path,
                left=Inches(7.45), top=Inches(0.9),
                width=Inches(5.5),
            )
        except Exception as e:
            logger.error(f"Failed to insert image on slide {slide_num}: {e}")
            # Gracefully fall back — right column stays empty rather than crashing

        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    def _add_closing_slide(self, prs, slide_info, slide_num, total):
        """Closing slide with gold accent — visually distinct from content slides."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = prs.slide_width, prs.slide_height

        self._set_bg(slide, self.COLOR_BG_ALT)

        # Gold bottom bar
        self._add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), self.COLOR_ACCENT2)

        # Gold left bar
        self._add_rect(slide, 0, 0, Inches(0.18), H, self.COLOR_ACCENT2)

        # Large decorative circle (background element)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(6.5), Inches(0.5), Inches(6.5), Inches(6.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0x25, 0x32, 0x75)
        circle.line.fill.background()

        # Title
        title_text = slide_info.get("title", "Thank You")
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(11.5), Inches(2.0))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(54)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_ACCENT2

        # Divider
        self._add_rect(slide, Inches(0.55), Inches(3.55), Inches(6.0), Inches(0.04), self.COLOR_ACCENT2)

        # Closing points
        content = self._coerce_list(slide_info.get("content", []))
        if content:
            tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(3.8), Inches(11.5), Inches(3.2))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for i, point in enumerate(content):
                para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                para.space_before = Pt(6)
                run_txt = para.add_run()
                run_txt.text = point
                run_txt.font.name  = self.FONT_BODY
                run_txt.font.size  = Pt(18)
                run_txt.font.color.rgb = self.COLOR_BODY

        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    # ──────────────────────────────────────────────────────────────────
    # Helpers
    # ──────────────────────────────────────────────────────────────────

    def _set_bg(self, slide, color: RGBColor):
        bg   = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_rect(self, slide, left, top, width, height, color: RGBColor):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_slide_number(self, slide, prs, slide_num: int, total: int):
        W = prs.slide_width
        H = prs.slide_height
        tb = slide.shapes.add_textbox(W - Inches(1.3), H - Inches(0.45), Inches(1.1), Inches(0.35))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{slide_num} / {total}"
        run.font.name  = self.FONT_BODY
        run.font.size  = Pt(10)
        run.font.color.rgb = self.COLOR_MUTED

    def _add_speaker_notes(self, slide, slide_info: dict):
        """Adds speaker notes and Creative Director visual instructions."""
        notes_text = slide_info.get("notes", "")
        
        # Append Visual Design Block if provided by the Architect
        visual = slide_info.get("visual")
        if visual and isinstance(visual, dict):
            visual_block = (
                f"\n\n🎨 VISUAL DESIGN (Creative Director)\n"
                f"--------------------------------------\n"
                f"IMAGE CONCEPT: {visual.get('concept', 'N/A')}\n\n"
                f"AI IMAGE PROMPT: {visual.get('prompt', 'N/A')}\n"
                f"STYLE: Clean 3D Isometric\n"
            )
            notes_text += visual_block

        if notes_text.strip():
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text.strip()

    @staticmethod
    def _coerce_list(content) -> list:
        if isinstance(content, str):
            return [content] if content else []
        return list(content) if content else []

    @staticmethod
    def _unique_path(path: str) -> str:
        if not os.path.exists(path):
            return path
        base, ext = os.path.splitext(path)
        counter = 1
        while os.path.exists(f"{base}_{counter}{ext}"):
            counter += 1
        return f"{base}_{counter}{ext}"